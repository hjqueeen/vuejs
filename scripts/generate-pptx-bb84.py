"""
BB84 Protocol — English Presentation Script → PowerPoint (.pptx)
Output: src/data/BB84_Protocol_en.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# ── palette ───────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x1B, 0x2A)
ACCENT  = RGBColor(0x5B, 0xCE, 0xFA)
VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
MUTED   = RGBColor(0x9C, 0xB4, 0xCC)
PANEL   = RGBColor(0x13, 0x27, 0x3D)

W = Inches(13.33)
H = Inches(7.5)

# ── slide data ────────────────────────────────────────────────────────────────
# (kind, title, bullets, speaker_notes)
SLIDES = [
    ("cover",
     "BB84 Protocol\nImplementation with Qiskit",
     ["Hyejin Kim"],
     ""),

    ("overview",
     "Agenda",
     [
         "① Basics of Quantum Computing",
         "② Limitations of Classical Cryptography",
         "③ Quantum Cryptography & QKD",
         "④ BB84 Protocol",
         "⑤ Hands-on Implementation with Qiskit",
     ],
     "Quick overview of what we'll cover today."),

    ("section", "Part 1\nQuantum Computing", [], ""),

    ("content",
     "What Is Quantum Computing?",
     [
         "Uses quantum phenomena to process data massively in parallel",
         "Computation speed: exponentially faster than classical computers",
         "Google's quantum supremacy — solved a 10.7-trillion-year problem in 5 min",
         "Driving huge investments from companies worldwide",
     ],
     "Hello, I'm Hyejin Kim, and today I'll be talking about implementing the BB84 protocol using Qiskit. In this presentation, we'll cover some basic concepts of quantum computing, the limitations of conventional cryptography, and quantum cryptography — including the BB84 protocol. After that, we'll walk through a hands-on implementation of BB84 with Qiskit. For the practical part, you'll need basic Python knowledge and familiarity with Google Colab.\n\nQuantum computing is a technology that uses quantum phenomena to process data in parallel on a massive scale, which makes computation extremely fast. A recent headline example is Google's quantum supremacy experiment — news reports said they solved a problem that would take about 10.7 trillion years on a classical computer in just five minutes. Because of this speed, quantum computing is expected to be useful across many fields, and companies around the world are investing heavily in its development."),

    ("section", "Part 2\nClassical Cryptography & Its Limits", [], ""),

    ("content",
     "How Classical Cryptography Works",
     [
         "Asymmetric public-key encryption — widely used today",
         "Hard mathematical problem → used as the public key",
         "Only the holder of the private key can decrypt",
         "Security = computational complexity",
     ],
     "Before we get into quantum cryptography, let me briefly explain how conventional cryptography works. The asymmetric public-key schemes widely used today rely on mathematical problems that are very hard to solve. You encrypt information with a public key, and only someone with the private key can decrypt it — so security is based on computational complexity."),

    ("content",
     "RSA & Integer Factorization",
     [
         "RSA uses the integer factorization problem",
         "Time to solve grows exponentially with problem size",
         "Large enough key → practically impossible for eavesdroppers to break",
     ],
     "Take RSA as an example: it uses the integer factorization problem. As the size of the problem grows, the time needed to solve it increases exponentially. So if sender and receiver use a large enough factorization problem as their public key, an eavesdropper who intercepts the ciphertext can't realistically decrypt it in practice."),

    ("content",
     "The Quantum Threat",
     [
         "Classical crypto relies on 'decryption takes too long'",
         "Quantum computers can solve hard math problems almost instantly",
         "Existing cryptographic systems can no longer be guaranteed safe",
     ],
     "Traditional cryptography stayed secure because decryption takes a long time — but as I mentioned, quantum computers can solve these mathematical hard problems in a flash. As quantum computing advances, we can no longer assume that our existing crypto systems are safe."),

    ("section", "Part 3\nQuantum Cryptography & QKD", [], ""),

    ("content",
     "Quantum Cryptography as the Solution",
     [
         "Transmit the encryption key using the laws of quantum mechanics",
         "Eavesdropping becomes physically impossible",
         "Security grounded in physical principles — not computational difficulty",
         "Safe regardless of quantum computing power",
     ],
     "So how do we address this? One answer is quantum cryptography. When you transmit the encryption key, quantum cryptography uses the laws of quantum mechanics to make eavesdropping impossible. Where classical crypto counted on slow decryption for safety, quantum cryptography uses physical principles so you maintain a consistent level of security regardless of how powerful quantum computers become."),

    ("content",
     "Three Quantum Properties",
     [
         "Superposition — qubit holds 0 and 1 simultaneously; key values stay hidden",
         "Measurement Uncertainty — measuring collapses the state; eavesdropping detectable",
         "No-Cloning Theorem — arbitrary quantum information cannot be perfectly copied",
     ],
     "Now let me explain the quantum properties that quantum cryptography relies on. Quantum systems have three key features: superposition, measurement uncertainty, and the no-cloning theorem. Superposition means a qubit can be in both 0 and 1 at once — we use that to send keys in overlapping states so their values stay hidden. Uncertainty means that before you measure a quantum state it can be in superposition, but once you measure, the state collapses — and that lets us detect eavesdropping during key exchange. No-cloning means you cannot perfectly copy arbitrary quantum information — so an eavesdropper can't measure, copy, and forward the same quantum state to hide what they did."),

    ("content",
     "QKD — Quantum Key Distribution",
     [
         "Distributes encryption keys securely in real time using quantum properties",
         "Also known as QKD (Quantum Key Distribution)",
         "BB84 is one of the most representative QKD protocols",
         "Uses four quantum states for key distribution",
     ],
     "Using these properties, quantum cryptography builds a system where sender and receiver can safely distribute encryption keys in real time. This is known as QKD — Quantum Key Distribution. Next, I'll walk you through BB84, one of the most representative QKD protocols. BB84 performs quantum key distribution using four quantum states. Before we look at the protocol steps, there are a few concepts we should cover — I'll keep this simple so it's not too overwhelming."),

    ("section", "Part 4\nBB84 Protocol", [], ""),

    ("content",
     "Qubit & Two Bases",
     [
         "Qubit — basic unit; probabilistic superposition of 0 and 1",
         "Rectilinear basis:  0 = horizontal pol.   |  1 = vertical pol.",
         "Diagonal basis:     0 = diagonal pol.     |  1 = anti-diagonal pol.",
         "Applying different polarization filters changes how the state is represented",
     ],
     "In quantum computing, the basic unit of information is the quantum bit, or qubit. Unlike a classical bit that is fixed at 0 or 1, a qubit can exist in a superposition where it probabilistically holds both 0 and 1 at the same time. A basis is a set of orthogonal vectors used to represent quantum states. BB84 uses two bases — the rectilinear basis and the diagonal basis — and when a quantum state passes through different polarization filters, its representation changes, as shown in the diagram. In the rectilinear basis, 0 maps to horizontal polarization and 1 to vertical. In the diagonal basis, 0 maps to diagonal polarization and 1 to anti-diagonal."),

    ("content",
     "BB84 Protocol Steps",
     [
         "① Alice generates random bits and randomly picks a basis for each",
         "② Alice encodes each bit using the chosen basis",
         "③ Bob independently picks a random basis and measures each qubit",
         "④ Alice & Bob compare bases publicly — keep only matching bits (key sifting)",
         "⑤ Keys match = no eavesdropping  /  Keys differ = eavesdropping detected",
     ],
     "Now let's go through how BB84 works. In cryptography we often name the sender Alice, the receiver Bob, and the eavesdropper Eve — I'll use those names here. First, Alice generates random bits and randomly chooses a basis for each one. Then she encodes each bit in the chosen basis — for example, bit 0 with the rectilinear basis becomes horizontal polarization, bit 1 with rectilinear becomes vertical, and with the diagonal basis, 0 becomes diagonal and 1 becomes anti-diagonal. Bob also picks a random basis for each incoming qubit. After transmission, Alice and Bob publicly compare their basis choices and keep only the bits where they matched — those become the shared secret key. If their keys agree, they assume the channel was secure; if not, they treat it as possible eavesdropping. That's the core idea of BB84."),

    ("section", "Part 5\nImplementation with Qiskit", [], ""),

    ("content",
     "What Is Qiskit?",
     [
         "Open-source quantum computing framework by IBM",
         "Access real quantum computers via the cloud & design circuits",
         "Lab will be run in Google Colab (code pre-written)",
         "Requirements: basic Python + basic Google Colab",
     ],
     "Now let's implement BB84 using Qiskit. Qiskit is an open-source quantum computing framework designed to access quantum computers through the cloud and to design quantum circuits. We'll run the lab in Google Colab, and the code is already prepared. First we'll install and import the packages Qiskit needs. Before the full BB84 implementation, let's look at a simple case: Alice sends one bit to Bob and Bob receives it."),

    ("content",
     "Hadamard Gate — Changing the Basis",
     [
         "Qubit 0 through diagonal basis -> encoded as diagonal polarization",
         "Hadamard gate: creates superposition, switches measurement basis",
         "Rectilinear basis -> no gate needed",
         "Diagonal basis   -> apply Hadamard gate",
     ],
     "When qubit 0 passes through the diagonal basis, it gets encoded as diagonal polarization. We can implement the diagonal basis using the Hadamard gate. The Hadamard gate creates superposition and is used to change the measurement basis — when you want the rectilinear basis you leave the qubit alone, and when you need the diagonal basis you apply a Hadamard gate."),

    ("content",
     "Demo 1 — No Eavesdropping",
     [
         "Create a 1-qubit circuit; initial state |0>",
         "Alice applies Hadamard -> switches to diagonal basis",
         "Bob applies Hadamard and measures in the same basis",
         "Result: 100% probability of measuring 0 (original state preserved)",
     ],
     "Let's look at the code. We create a one-qubit quantum circuit. Right after initialization the qubit starts in |0>, and we apply a Hadamard gate to switch it to the diagonal basis — so 0 becomes diagonally polarized. The barrier in the diagram is just for visual separation; it doesn't change the physics. Bob then applies another Hadamard gate and measures in the same diagonal basis to decode. When we run the circuit, we see the original state Alice sent — 0 — with 100% probability."),

    ("content",
     "Demo 2 — With Eavesdropping (Eve)",
     [
         "Eve intercepts and measures the qubit mid-transmission",
         "Measurement collapses the quantum state",
         "Bob can no longer recover the original value",
         "Result: 0 and 1 appear with ~50% probability each",
     ],
     "Next, the eavesdropping scenario. As I said, because of measurement uncertainty, once Eve measures the qubit it collapses and Bob can no longer recover the original state. The code is the same as before, except we insert Eve's measurement in the middle. After that intervention, the result is no longer 100% zero — we get 0 and 1 with about 50% each."),

    ("content",
     "Full BB84 Lab — No Eavesdropping",
     [
         "NumPy random seed -> generate 35 bits and 35 bases",
         "encode_message(): build quantum circuits per bit & basis",
         "decode_message(): Bob picks random bases and measures",
         "generate_encryption_key(): keep only same-basis bits",
         "Keys match -> Transmission successful",
     ],
     "Now let's move on to the BB84 lab. The flow matches the slides, so I'll go straight into the implementation. First, the case with no eavesdropping: when Alice and Bob chose the same basis, Bob's decoded bits should match Alice's bits perfectly — and we use those matching cases to build the key.\n\nLet's verify that. Using NumPy's random seed, we'll randomly pick 35 bits and 35 bases. Think of basis 0 as rectilinear and basis 1 as diagonal. You can see the bits and bases were selected correctly.\n\nNext we encode according to the chosen bases. For each of the n bits we build a quantum circuit and set the state based on the bit and basis. If the basis is rectilinear and the bit is 0, we pass it through unchanged. If it's 1, we flip it with an X gate. For the diagonal basis with bit 0, we apply Hadamard; for bit 1 we X first, then Hadamard. We concatenate these encoded circuits in the encode_message function.\n\nBob randomly chooses bases for each qubit, then decodes. For n qubits, if Bob's basis is rectilinear he measures directly; if it's diagonal he applies Hadamard first, then measures.\n\nAlice and Bob now publicly share their basis choices and add matching cases to the encryption key. The keys they generate are identical — run it once and you get a successful transmission message. When both keys match, we can conclude there was no eavesdropping."),

    ("content",
     "Full BB84 Lab — With Eve",
     [
         "Eve randomly picks a basis and decodes Alice's message mid-transit",
         "Eve's measurement collapses the quantum state",
         "Bob's decoded bits no longer match Alice's",
         "Even same-basis bits disagree -> eavesdropping detected",
     ],
     "As a bonus, here's what happens when someone eavesdrops. The code is mostly the same, but we added a step where Eve randomly picks a basis and decodes Alice's message. When Eve measures once, the quantum state collapses, so Bob can no longer keep the same result as Alice's original bit. That means even the bits where Alice and Bob chose the same basis no longer match perfectly — and that's how we detect eavesdropping."),

    ("closing",
     "Thank You",
     [
         "References",
         "   Several papers on quantum cryptography",
         "   Qiskit Textbook — github.com/Qiskit/textbook",
         "",
         "BB84 Protocol Implementation with Qiskit",
         "Hyejin Kim",
     ],
     "For references, I consulted several papers on quantum cryptography and the Qiskit textbook on GitHub for the implementation. That wraps up our look at implementing the BB84 protocol with Qiskit. Thank you."),
]

# ── helpers ───────────────────────────────────────────────────────────────────
def rect(slide, l, t, w, h, color):
    shp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def txbox(slide, text, l, t, w, h, size=18, bold=False,
          color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box

def notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text

def bg(slide):
    rect(slide, 0, 0, W, H, BG)

# ── render each slide ─────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

for kind, title, bullets, speaker in SLIDES:
    sl = prs.slides.add_slide(blank)
    bg(sl)

    # ── COVER ────────────────────────────────────────────────────────────────
    if kind == "cover":
        rect(sl, 0, 0,           W,          Inches(0.09), ACCENT)
        rect(sl, 0, H - Inches(0.09), W,     Inches(0.09), VIOLET)
        rect(sl, Inches(0.6), Inches(2.1), Inches(12.1), Inches(3.4), PANEL)

        txbox(sl, "QUANTUM CRYPTOGRAPHY  |  QISKIT",
              Inches(0.85), Inches(2.3), Inches(11), Inches(0.5),
              size=11, bold=True, color=ACCENT)

        parts = title.split("\n")
        txbox(sl, parts[0],
              Inches(0.85), Inches(2.75), Inches(11), Inches(1.1),
              size=44, bold=True, color=WHITE)
        if len(parts) > 1:
            txbox(sl, parts[1],
                  Inches(0.85), Inches(3.75), Inches(11), Inches(0.7),
                  size=26, color=ACCENT)

        txbox(sl, bullets[0],
              Inches(0.85), Inches(4.65), Inches(8), Inches(0.5),
              size=16, color=MUTED)

    # ── SECTION divider ──────────────────────────────────────────────────────
    elif kind == "section":
        rect(sl, 0, 0,           Inches(0.2), H,          ACCENT)
        rect(sl, 0, H - Inches(0.09), W,     Inches(0.09), VIOLET)

        parts = title.split("\n")
        txbox(sl, parts[0],
              Inches(0.55), Inches(2.85), Inches(12), Inches(0.65),
              size=20, bold=True, color=ACCENT)
        if len(parts) > 1:
            txbox(sl, parts[1],
                  Inches(0.55), Inches(3.45), Inches(12), Inches(1.1),
                  size=40, bold=True, color=WHITE)

    # ── OVERVIEW ─────────────────────────────────────────────────────────────
    elif kind == "overview":
        rect(sl, 0, 0,           W, Inches(0.09), ACCENT)
        rect(sl, 0, H - Inches(0.09), W, Inches(0.09), VIOLET)
        rect(sl, 0, Inches(0.09), W, Inches(1.1), PANEL)

        txbox(sl, title,
              Inches(0.4), Inches(0.16), Inches(12.5), Inches(0.9),
              size=28, bold=True, color=WHITE)

        for i, b in enumerate(bullets):
            top = Inches(1.45) + Inches(i * 0.72)
            rect(sl, Inches(0.45), top + Inches(0.17),
                 Inches(0.13), Inches(0.13), ACCENT)
            txbox(sl, b, Inches(0.72), top, Inches(12), Inches(0.6),
                  size=20, color=WHITE)

    # ── CONTENT ──────────────────────────────────────────────────────────────
    elif kind == "content":
        rect(sl, 0, 0,           W, Inches(0.09), ACCENT)
        rect(sl, 0, H - Inches(0.09), W, Inches(0.09), VIOLET)
        rect(sl, 0, Inches(0.09), W, Inches(1.1), PANEL)

        txbox(sl, title,
              Inches(0.4), Inches(0.16), Inches(12.5), Inches(0.9),
              size=26, bold=True, color=WHITE)

        step = Inches(0.72) if len(bullets) <= 5 else Inches(0.62)
        for i, b in enumerate(bullets):
            if not b:
                continue
            top = Inches(1.4) + i * step
            rect(sl, Inches(0.45), top + Inches(0.17),
                 Inches(0.13), Inches(0.13), ACCENT)
            txbox(sl, b, Inches(0.72), top, Inches(12.2), Inches(0.58),
                  size=18, color=WHITE)

    # ── CLOSING ──────────────────────────────────────────────────────────────
    elif kind == "closing":
        rect(sl, 0, 0,           W, Inches(0.09), ACCENT)
        rect(sl, 0, H - Inches(0.09), W, Inches(0.09), VIOLET)
        rect(sl, Inches(0.6), Inches(1.7), Inches(12.1), Inches(4.5), PANEL)

        txbox(sl, title,
              0, Inches(0.85), W, Inches(0.8),
              size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        y = Inches(2.0)
        for line in bullets:
            indent = line.startswith("   ")
            c = MUTED if indent else ACCENT
            sz = 14 if indent else 17
            txbox(sl, line.strip(), Inches(1.4), y, Inches(11), Inches(0.46),
                  size=sz, color=c)
            y += Inches(0.44)

    notes(sl, speaker)

# ── save ──────────────────────────────────────────────────────────────────────
out = os.path.normpath(
    os.path.join(os.path.dirname(__file__), "../src/data/BB84_Protocol_en.pptx"))
prs.save(out)
print(f"Saved -> {out}  ({os.path.getsize(out)//1024} KB)")
